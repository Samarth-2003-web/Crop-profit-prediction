"""
================================================================================
AGRICULTURAL PROFIT PREDICTOR - SACAIM 2025 PRESENTATION GENERATOR
================================================================================
This script creates a complete PowerPoint presentation for the conference.
Run: python create_presentation.py
================================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

print("="*80)
print(" " * 20 + "GENERATING SACAIM 2025 PRESENTATION")
print("="*80)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list, bg_color=RGBColor(255, 255, 255)):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_shape.line.color.rgb = RGBColor(41, 128, 185)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_before = Pt(12)
        p.level = 0

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(41, 128, 185)
    title_shape.line.color.rgb = RGBColor(41, 128, 185)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    for i, item in enumerate(left_content):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_before = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    for i, item in enumerate(right_content):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.space_before = Pt(10)

# ============================================================================
# SLIDE 1: TITLE SLIDE
# ============================================================================
print("\nCreating Slide 1: Title Slide...")
add_title_slide(prs, 
    "Agricultural Profit Predictor",
    "Machine Learning for Indian Farmers\nSACAIM 2025 Conference"
)

# ============================================================================
# SLIDE 2: INTRODUCTION
# ============================================================================
print("Creating Slide 2: Introduction...")
add_content_slide(prs, "Introduction", [
    "• Agriculture is the backbone of India's economy",
    "• Farmers face challenges in crop selection and profit prediction",
    "• Traditional methods lack data-driven insights",
    "• Machine Learning can provide accurate predictions",
    "• Our solution: Web-based profit predictor with 92.37% accuracy",
    "• Helps farmers make informed decisions about crop selection"
])

# ============================================================================
# SLIDE 3: PROBLEM STATEMENT
# ============================================================================
print("Creating Slide 3: Problem Statement...")
add_content_slide(prs, "Problem Statement", [
    "CHALLENGE:",
    "  - Farmers struggle to predict crop profitability",
    "  - Limited access to historical yield and price data",
    "  - Complex factors affect crop performance (soil, weather, location)",
    "",
    "IMPACT:",
    "  - Poor crop selection leads to financial losses",
    "  - Uncertainty in planning and resource allocation",
    "",
    "OUR SOLUTION:",
    "  - ML-based system predicting yield, price, and profit",
    "  - Simple 8-input interface for farmers",
    "  - Real-time predictions with high accuracy"
])

# ============================================================================
# SLIDE 4: DATA COLLECTION
# ============================================================================
print("Creating Slide 4: Data Collection...")
add_two_column_slide(prs, "Data Collection", 
    [
        "YIELD DATA:",
        "• Source: data.gov.in",
        "• Region: Karnataka, India",
        "• Records: 19,171 observations",
        "• Time Period: 2015-2025",
        "• Features: 14 attributes",
        "• Crops: 10 major crops",
        "• Districts: 30 districts",
        "• Seasons: Kharif, Rabi, Summer"
    ],
    [
        "PRICE DATA:",
        "• Source: Krishimaratavahini Portal",
        "• Government of Karnataka portal",
        "• Records: 1,320 price observations",
        "• Time Period: 2015-2025",
        "• Monthly price trends",
        "• APMC market rates",
        "• Inflation-adjusted prices",
        "• Base prices for 10 crops"
    ]
)

# ============================================================================
# SLIDE 5: DATA PREPROCESSING
# ============================================================================
print("Creating Slide 5: Data Preprocessing...")
add_content_slide(prs, "Data Preprocessing", [
    "1. DATA CLEANING:",
    "   • Removed duplicate records",
    "   • Handled missing values",
    "   • Fixed inconsistent data formats",
    "",
    "2. DATA NORMALIZATION:",
    "   • Standardized numerical features",
    "   • Scaled features to common range",
    "",
    "3. ENCODING:",
    "   • Label Encoding for Season (Kharif, Rabi, Summer)",
    "   • Label Encoding for Crop Type (10 crops)",
    "   • Label Encoding for District (30 districts)"
])

# ============================================================================
# SLIDE 6: FEATURE ENGINEERING
# ============================================================================
print("Creating Slide 6: Feature Engineering...")
add_two_column_slide(prs, "Feature Engineering",
    [
        "USER INPUTS (8):",
        "1. Season",
        "2. Crop Type",
        "3. District",
        "4. Year",
        "5. Area (acres)",
        "6. Previous Production",
        "7. Previous Area",
        "8. Previous Cost"
    ],
    [
        "ENGINEERED FEATURES (15):",
        "• season_enc, crop_enc, dist_enc",
        "• year, area",
        "• soil_ph, soil_N, soil_P, soil_K",
        "• organic_matter, irrigation_type",
        "• rainfall_mm, avg_temperature",
        "• soil_fertility (N×P×K)^(1/3)",
        "• water_temp_interaction",
        "",
        "DEFAULT PARAMETERS:",
        "• Soil & weather data from research"
    ]
)

# ============================================================================
# SLIDE 7: MODEL SELECTION
# ============================================================================
print("Creating Slide 7: Model Selection...")
add_content_slide(prs, "Model Selection & Training", [
    "THREE ALGORITHMS COMPARED:",
    "",
    "1. RANDOM FOREST REGRESSOR",
    "   • 400 trees, max_depth=25",
    "   • Good for handling non-linear relationships",
    "",
    "2. GRADIENT BOOSTING REGRESSOR",
    "   • 300 estimators, learning_rate=0.05",
    "   • Sequential learning from errors",
    "",
    "3. XGBOOST REGRESSOR (WINNER)",
    "   • 400 estimators, learning_rate=0.05",
    "   • Parallel processing, regularization",
    "   • Best performance on validation set"
])

# ============================================================================
# SLIDE 8: MODEL TRAINING
# ============================================================================
print("Creating Slide 8: Model Training Process...")
add_content_slide(prs, "Model Training Process", [
    "TRAINING SETUP:",
    "  • Train-Test Split: 80% - 20%",
    "  • Training Samples: 15,336",
    "  • Test Samples: 3,835",
    "  • Cross-validation: 5-fold",
    "",
    "HYPERPARAMETER TUNING:",
    "  • Grid Search for optimal parameters",
    "  • Max Depth: 8-25",
    "  • Learning Rate: 0.01-0.1",
    "  • N_estimators: 100-500",
    "",
    "TRAINING TIME:",
    "  • XGBoost: ~45 seconds",
    "  • Random Forest: ~38 seconds",
    "  • Gradient Boosting: ~52 seconds"
])

# ============================================================================
# SLIDE 9: RESULTS & ACCURACY
# ============================================================================
print("Creating Slide 9: Results & Accuracy...")
add_content_slide(prs, "Results & Accuracy", [
    "MODEL PERFORMANCE COMPARISON:",
    "",
    "XGBOOST (BEST MODEL):",
    "  • Yield Prediction: R² = 0.9240 (92.40%)",
    "  • Price Prediction: R² = 0.9703 (97.03%)",
    "",
    "RANDOM FOREST:",
    "  • Yield Prediction: R² = 0.9237 (92.37%)",
    "  • Price Prediction: R² = 0.9701 (97.01%)",
    "",
    "GRADIENT BOOSTING:",
    "  • Yield Prediction: R² = 0.9185 (91.85%)",
    "  • Price Prediction: R² = 0.9658 (96.58%)",
    "",
    "CONFUSION MATRIX: 85.71% category accuracy (Low/Medium/High)"
])

# ============================================================================
# SLIDE 10: BEST MODEL DETAILS
# ============================================================================
print("Creating Slide 10: Best Model...")
add_content_slide(prs, "Best Model: XGBoost", [
    "WHY XGBOOST PERFORMS BEST:",
    "",
    "1. REGULARIZATION:",
    "   • Built-in L1 and L2 regularization prevents overfitting",
    "",
    "2. PARALLEL PROCESSING:",
    "   • Faster training with multi-core support",
    "",
    "3. HANDLING MISSING DATA:",
    "   • Native support for sparse data",
    "",
    "4. FEATURE IMPORTANCE:",
    "   • Identifies key predictors (area, soil_fertility, crop_type)",
    "",
    "5. ROBUSTNESS:",
    "   • Performs well across different yield ranges",
    "   • Minimal misclassification between Low and High yields"
])

# ============================================================================
# SLIDE 11: WEB APPLICATION
# ============================================================================
print("Creating Slide 11: Web Application...")
add_two_column_slide(prs, "Web Application Deployment",
    [
        "TECHNOLOGY STACK:",
        "• Backend: Flask (Python)",
        "• Frontend: HTML5, CSS3, JavaScript",
        "• ML Library: XGBoost, Scikit-learn",
        "• Deployment: Render Cloud",
        "• Response Time: <1 second",
        "",
        "FEATURES:",
        "• Simple 8-input form",
        "• Real-time predictions",
        "• Mobile-responsive design",
        "• No registration required"
    ],
    [
        "OUTPUT DASHBOARD:",
        "• Predicted Yield (quintals/acre)",
        "• Market Price (₹/quintal)",
        "• Total Revenue (₹)",
        "• Total Cost (₹)",
        "• Net Profit (₹)",
        "• Return on Investment (%)",
        "",
        "FARMER BENEFITS:",
        "• Informed crop selection",
        "• Risk assessment",
        "• Profit maximization",
        "• Easy-to-use interface"
    ]
)

# ============================================================================
# SLIDE 12: CONFUSION MATRIX ANALYSIS
# ============================================================================
print("Creating Slide 12: Confusion Matrix...")
add_content_slide(prs, "Confusion Matrix Analysis", [
    "YIELD CATEGORIZATION:",
    "  • Low Yield: < 0.0752 tonnes/hectare (33%)",
    "  • Medium Yield: 0.0752 - 0.1177 tonnes/hectare (33%)",
    "  • High Yield: > 0.1177 tonnes/hectare (34%)",
    "",
    "CLASSIFICATION ACCURACY:",
    "  • Low Yield: 86.15% accuracy (1,101/1,278 correct)",
    "  • Medium Yield: 80.14% accuracy (1,025/1,279 correct)",
    "  • High Yield: 90.85% accuracy (1,161/1,278 correct)",
    "  • Overall: 85.71% category classification accuracy",
    "",
    "KEY INSIGHT:",
    "  • No extreme errors (Low-High confusion = 0)",
    "  • Model reliably distinguishes good crops from bad ones"
])

# ============================================================================
# SLIDE 13: CONCLUSION
# ============================================================================
print("Creating Slide 13: Conclusion...")
add_content_slide(prs, "Conclusion", [
    "ACHIEVEMENTS:",
    "  • Developed ML-based profit predictor with 92.40% accuracy",
    "  • Deployed user-friendly web application for farmers",
    "  • Processed 19,171 historical records (2015-2025)",
    "  • Achieved 97.03% price prediction accuracy",
    "",
    "IMPACT:",
    "  • Helps farmers make data-driven decisions",
    "  • Reduces financial risks in crop selection",
    "  • Accessible via any device with internet",
    "",
    "INNOVATION:",
    "  • Only 8 simple inputs required from farmers",
    "  • Automated feature engineering (15 features)",
    "  • Real-time predictions (<1 second)",
    "  • First of its kind for Karnataka agriculture"
])

# ============================================================================
# SLIDE 14: FUTURE WORK
# ============================================================================
print("Creating Slide 14: Future Work...")
add_content_slide(prs, "Future Work & Enhancements", [
    "PLANNED ENHANCEMENTS:",
    "",
    "1. API INTEGRATION:",
    "   • Real-time weather data from meteorological APIs",
    "   • Live soil data from agricultural databases",
    "",
    "2. EXPANDED COVERAGE:",
    "   • More states across India",
    "   • Additional crops and varieties",
    "",
    "3. MOBILE APPLICATION:",
    "   • Native Android/iOS apps",
    "   • Offline prediction capability",
    "",
    "4. MULTILINGUAL SUPPORT:",
    "   • Kannada, Hindi, Telugu, Tamil interfaces",
    "",
    "5. ADVANCED FEATURES:",
    "   • Disease prediction and prevention",
    "   • Optimal planting time recommendations",
    "   • Fertilizer optimization suggestions"
])

# ============================================================================
# SLIDE 15: THANK YOU
# ============================================================================
print("Creating Slide 15: Thank You Slide...")
add_title_slide(prs,
    "Thank You!",
    "Questions & Discussion\n\nSACAIM 2025 Conference"
)

# ============================================================================
# SAVE PRESENTATION
# ============================================================================
output_file = 'Agricultural_Profit_Predictor_SACAIM_2025.pptx'
prs.save(output_file)

print("\n" + "="*80)
print("✅ PRESENTATION CREATED SUCCESSFULLY!")
print("="*80)
print(f"📁 File: {output_file}")
print(f"📊 Total Slides: 15")
print(f"📏 Size: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
print(f"🎨 Professional design with consistent formatting")
print("\nSlide Breakdown:")
print("  1. Title Slide")
print("  2. Introduction")
print("  3. Problem Statement")
print("  4. Data Collection")
print("  5. Data Preprocessing")
print("  6. Feature Engineering")
print("  7. Model Selection")
print("  8. Model Training")
print("  9. Results & Accuracy")
print("  10. Best Model (XGBoost)")
print("  11. Web Application")
print("  12. Confusion Matrix")
print("  13. Conclusion")
print("  14. Future Work")
print("  15. Thank You")
print("\n" + "="*80)
print("📌 Next Step: Open the file in PowerPoint to add images/charts!")
print("="*80)
